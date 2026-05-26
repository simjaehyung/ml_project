"""
Hotel No-Show DSS v12.html → PPTX 변환 스크립트
Reveal.js HTML 슬라이드를 파싱해 python-pptx로 재구성
"""
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
import io

# ── 색상 팔레트 ─────────────────────────────────────────────
BRAND     = RGBColor(0x1E, 0x3A, 0x8A)   # 진파랑
ACCENT    = RGBColor(0xD9, 0x77, 0x06)   # 골드
INK       = RGBColor(0x0F, 0x17, 0x2A)   # 거의 검정
INK2      = RGBColor(0x47, 0x55, 0x69)   # 회색
INK3      = RGBColor(0x94, 0xA3, 0xB8)   # 연회색
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_S    = RGBColor(0xEF, 0xF6, 0xFF)   # 연파랑 배경
GOLD_S    = RGBColor(0xFF, 0xFB, 0xEB)   # 연골드 배경
GREEN     = RGBColor(0x05, 0x96, 0x69)
RED       = RGBColor(0xDC, 0x26, 0x26)
BD        = RGBColor(0xE2, 0xE8, 0xF0)   # 테두리

# ── 슬라이드 크기 (16:9) ────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_header(slide, slide_num: int, total: int = 14, dark=False):
    """상단 브랜드 바"""
    brand_color = RGBColor(0xFF,0xFF,0xFF) if dark else INK3
    add_text_box(slide, "HOTEL NO-SHOW  ·  DSS",
                 Inches(0.5), Inches(0.18), Inches(9), Inches(0.3),
                 font_size=9, bold=True, color=brand_color)
    add_text_box(slide, f"{slide_num:02d} / {total}",
                 Inches(11.5), Inches(0.18), Inches(1.5), Inches(0.3),
                 font_size=9, color=brand_color, align=PP_ALIGN.RIGHT)

def add_divider(slide, y=Inches(0.56), dark=False):
    from pptx.util import Pt as UPt
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(0.5), y, Inches(12.83), y
    )
    line.line.color.rgb = BD if not dark else RGBColor(0x2D,0x3A,0x55)
    line.line.width = Pt(0.5)

def label_tag(slide, text, left, top):
    """골드 소문자 레이블"""
    add_text_box(slide, text, left, top, Inches(8), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT)

def title_block(slide, main, sub=None, left=Inches(0.6), top=Inches(0.7)):
    """슬라이드 제목 (lbl + ttl)"""
    if sub:
        label_tag(slide, sub, left, top)
        top += Inches(0.32)
    add_text_box(slide, main, left, top, Inches(12.2), Inches(1.1),
                 font_size=34, bold=True, color=INK)

def bullet_item(tf, text, level=0, size=15, color=INK2, bold=False):
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    indent = "  " * level
    run = p.add_run()
    run.text = indent + ("• " if level == 0 else "– ") + text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold

def add_bullet_box(slide, items, left, top, width, height,
                   size=15, color=INK2, title=None, title_color=BRAND):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(13)
        run.font.bold  = True
        run.font.color.rgb = title_color
    for i, item in enumerate(items):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "• " + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color

def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None):
    rows_count = len(rows) + 1
    cols_count = len(headers)
    tbl = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    tbl.first_row = True

    # 컬럼 너비
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    # 헤더
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF8,0xFA,0xFC)
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = h
        run.font.size  = Pt(11)
        run.font.bold  = True
        run.font.color.rgb = INK3

    # 데이터
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.text = str(val)
            run.font.size  = Pt(13)
            run.font.color.rgb = INK2

    return tbl

def add_callout(slide, text, left, top, width, height, style="gold"):
    colors = {
        "gold": (ACCENT, GOLD_S),
        "blue": (BRAND,  BLUE_S),
        "green":(GREEN,  RGBColor(0xF0,0xFD,0xF4)),
    }
    bar_c, bg_c = colors.get(style, colors["gold"])
    # 배경 사각형
    bg = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_c
    bg.line.color.rgb = bar_c
    bg.line.width = Pt(0.5)
    # 텍스트
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05),
                                     width - Inches(0.15), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(14)
    run.font.color.rgb = INK

# ═══════════════════════════════════════════════════════════
# 슬라이드별 생성 함수
# ═══════════════════════════════════════════════════════════

def slide_01(prs):
    """타이틀 슬라이드"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 1)
    add_divider(sl)

    add_text_box(sl, "2026 · 데이터사이언스 프로젝트",
                 Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
                 font_size=12, bold=True, color=ACCENT)

    add_text_box(sl, "취소를 예측하고,",
                 Inches(0.6), Inches(1.5), Inches(10), Inches(0.8),
                 font_size=52, bold=True, color=INK)
    add_text_box(sl, "결정을 지원합니다.",
                 Inches(0.6), Inches(2.2), Inches(10), Inches(0.8),
                 font_size=52, bold=True, color=ACCENT)

    # 핵심 산출물 4개
    items = [
        ("Q1. 몇 개 오버부킹?",   "ML 예측 + 수학적 최적화",    Inches(0.6)),
        ("Q2. 재배치 시 얼마?",    "LLM 에이전트 시뮬레이션",    Inches(3.5)),
        ("LightGBM PR-AUC 0.82", "점유율 연동 임계값",         Inches(6.4)),
        ("1,240 runs",            "LLM 보상 협상 · Streamlit 앱", Inches(9.3)),
    ]
    for title, sub, x in items:
        add_text_box(sl, title, x, Inches(3.3), Inches(2.7), Inches(0.4),
                     font_size=14, bold=True, color=BRAND)
        add_text_box(sl, sub,   x, Inches(3.75), Inches(2.7), Inches(0.3),
                     font_size=11, color=INK2)

    add_text_box(sl, "심재형 · 이고은 · 김나리  |  한양대학교",
                 Inches(0.6), Inches(6.9), Inches(10), Inches(0.35),
                 font_size=11, color=INK3)

def slide_02(prs):
    """The Problem — 오버부킹 딜레마"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 2)
    add_divider(sl)

    label_tag(sl, "The Problem", Inches(0.6), Inches(0.68))
    add_text_box(sl, "오버부킹 딜레마 — 두 방향 모두 비용이 있습니다.",
                 Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                 font_size=30, bold=True, color=INK)

    add_text_box(sl, "아무것도 하지 않으면 빈방을 잃고, 과도하게 오버부킹하면 손님 불만을 삽니다.",
                 Inches(0.6), Inches(1.75), Inches(11), Inches(0.45),
                 font_size=16, color=INK2)

    # 카드 2개
    cards = [
        ("빈방 손실", "37%", "평균 취소율\nCity 41.7% · Resort 27.8%\n\n예약 10개 중 3~4개가\n체크인 전 사라집니다.\n빈방은 팔 수 없는 시간입니다.", RED),
        ("손님 재배치 보상", "€300+", "오버부킹 발생 시 보상 비용\n예약한 방이 없을 때\n다른 호텔로 모심\n보상금 + 대안 호텔 비용 + 평판", ACCENT),
    ]
    for i, (title, big, body, color) in enumerate(cards):
        x = Inches(0.6 + i * 6.3)
        bg = sl.shapes.add_shape(1, x, Inches(2.35), Inches(5.9), Inches(3.9))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF,0xF5,0xF5) if color==RED else GOLD_S
        bg.line.color.rgb = RGBColor(0xFE,0xCA,0xCA) if color==RED else RGBColor(0xFD,0xE6,0x8A)
        bg.line.width = Pt(1)

        add_text_box(sl, title, x+Inches(0.2), Inches(2.55), Inches(5.5), Inches(0.4),
                     font_size=13, bold=True, color=INK)
        add_text_box(sl, big,   x+Inches(0.2), Inches(3.0),  Inches(5.5), Inches(0.7),
                     font_size=44, bold=True, color=color)
        add_text_box(sl, body,  x+Inches(0.2), Inches(3.7),  Inches(5.5), Inches(2.3),
                     font_size=14, color=INK2)

def slide_03(prs):
    """우리의 접근 — 두 가지 질문"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 3)
    add_divider(sl)

    label_tag(sl, "우리의 접근", Inches(0.6), Inches(0.68))
    add_text_box(sl, "오버부킹은 두 가지 질문입니다.",
                 Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                 font_size=30, bold=True, color=INK)
    add_text_box(sl, "각 질문이 다른 도구를 요구합니다 — 수학과 시뮬레이션을 분리해서 씁니다.",
                 Inches(0.6), Inches(1.75), Inches(11), Inches(0.4),
                 font_size=15, color=INK2)

    cols = [
        ("Q1 · 수학", "몇 개 오버부킹?", [
            "LightGBM으로 각 예약의 취소 확률 계산",
            "점유율 연동 임계값으로 Flexi 풀 크기 결정",
            "재배치율 < 2% 제약 하에 RevPAR 최대화",
        ], BRAND, BLUE_S),
        ("Q2 · 시뮬레이션", "재배치 시 얼마?", [
            "취소가 예상보다 적으면 일부 손님 재배치 발생",
            "보상 협상 이력이 실제로 존재하지 않음",
            "LLM 에이전트로 합성 데이터 생성",
        ], ACCENT, GOLD_S),
    ]
    for i, (lbl, ttl, bullets, c, bg_c) in enumerate(cols):
        x = Inches(0.6 + i * 6.4)
        bg = sl.shapes.add_shape(1, x, Inches(2.35), Inches(6.1), Inches(4.1))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_c
        bg.line.color.rgb = c
        bg.line.width = Pt(0.75)

        add_text_box(sl, lbl, x+Inches(0.2), Inches(2.55), Inches(3.7), Inches(0.35),
                     font_size=11, bold=True, color=c)
        add_text_box(sl, ttl, x+Inches(0.2), Inches(2.9),  Inches(3.7), Inches(0.5),
                     font_size=20, bold=True, color=INK)
        tb = sl.shapes.add_textbox(x+Inches(0.2), Inches(3.5), Inches(3.7), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"{j+1}.  {b}"
            run.font.size  = Pt(13)
            run.font.color.rgb = INK2

def slide_04(prs):
    """Q1 — 취소 예측 기반 오버부킹"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 4)
    add_divider(sl)

    label_tag(sl, "Q1 · 몇 개 오버부킹?", Inches(0.6), Inches(0.68))
    add_text_box(sl, "취소 예측 기반으로 수학적으로 결정합니다.",
                 Inches(0.6), Inches(1.0), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)
    add_text_box(sl, "119,390건 예약 학습 → LightGBM 취소 확률 → 점유율 임계값 → RevPAR 최적화",
                 Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
                 font_size=14, color=INK2)

    # PR-AUC 테이블
    add_text_box(sl, "PR-AUC 모델 비교 · 테스트셋 취소율 38.7%",
                 Inches(0.6), Inches(2.25), Inches(6), Inches(0.35),
                 font_size=11, bold=True, color=INK3)
    add_table(sl,
        ["모델", "PR-AUC", "비고"],
        [
            ["LightGBM", "0.8189", "✓ 최종 선택 · Dummy 대비 +0.43"],
            ["XGBoost",  "0.810",  ""],
            ["Random Forest", "0.791", ""],
            ["Logistic Reg.", "0.752", ""],
            ["Dummy",    "0.387", "기준선"],
        ],
        Inches(0.6), Inches(2.65), Inches(5.8), Inches(2.5),
        col_widths=[Inches(1.8), Inches(1.3), Inches(2.7)],
    )

    # 점유율 임계값 테이블
    add_text_box(sl, "점유율 연동 임계값 전략",
                 Inches(7.0), Inches(2.25), Inches(5.7), Inches(0.35),
                 font_size=11, bold=True, color=INK3)
    add_table(sl,
        ["점유율", "임계값", "전략"],
        [
            ["< 60%",  "0.55", "적극"],
            ["60~80%", "0.60", "표준 ★"],
            ["> 80%",  "0.65", "신중"],
        ],
        Inches(7.0), Inches(2.65), Inches(5.7), Inches(1.6),
        col_widths=[Inches(1.5), Inches(1.5), Inches(2.7)],
    )

    # RevPAR 하이라이트
    add_text_box(sl, "+81%",
                 Inches(7.0), Inches(4.4), Inches(3), Inches(0.8),
                 font_size=48, bold=True, color=GREEN)
    add_text_box(sl, "RevPAR 개선 (시뮬레이션)",
                 Inches(7.0), Inches(5.2), Inches(5.7), Inches(0.35),
                 font_size=14, color=INK2)
    add_text_box(sl, "walk rate < 2%  ·  재배치 제약 충족",
                 Inches(7.0), Inches(5.6), Inches(5.7), Inches(0.35),
                 font_size=13, color=INK3)

def slide_05(prs):
    """핵심 개념 — Flexi 시스템"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 5)
    add_divider(sl)

    label_tag(sl, "핵심 개념 · Flexi 시스템", Inches(0.6), Inches(0.68))
    add_text_box(sl, "오버부킹을 수익으로 바꾸는 메커니즘",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)
    add_text_box(sl, "취소 위험 예약을 별도 풀에 모아 할인으로 유도 — 빈방을 줄이면서도 손님을 잃지 않습니다.",
                 Inches(0.6), Inches(1.75), Inches(11.5), Inches(0.4),
                 font_size=14, color=INK2)

    # 플로우 스텝
    steps = [
        ("1", "취소 확률 계산", "LightGBM → 예약마다 0~1 점수"),
        ("2", "임계값 판별",    "점유율 따라 0.55 / 0.60 / 0.65"),
        ("3a","Flexi 풀 배정", "확률 > 임계값\n할인 5~18% 제공"),
        ("3b","Standard 배정", "확률 ≤ 임계값\n일반 처리"),
    ]
    xs = [Inches(0.6), Inches(3.4), Inches(6.2), Inches(9.0)]
    for (num, ttl, body), x in zip(steps, xs):
        bg = sl.shapes.add_shape(1, x, Inches(2.35), Inches(2.7), Inches(3.4))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLUE_S if num in ("1","2","3b") else GOLD_S
        bg.line.color.rgb = BRAND if num in ("1","2","3b") else ACCENT
        bg.line.width = Pt(0.75)
        add_text_box(sl, num,  x+Inches(0.15), Inches(2.5),  Inches(0.5), Inches(0.4),
                     font_size=22, bold=True, color=BRAND if num in ("1","2","3b") else ACCENT)
        add_text_box(sl, ttl,  x+Inches(0.15), Inches(3.0),  Inches(2.4), Inches(0.45),
                     font_size=15, bold=True, color=INK)
        add_text_box(sl, body, x+Inches(0.15), Inches(3.5),  Inches(2.4), Inches(1.0),
                     font_size=12, color=INK2)

    # RevPAR 결과
    add_text_box(sl, "↑ RevPAR 상승 — 빈방 최소화 = 수익 최대화",
                 Inches(0.6), Inches(5.9), Inches(11), Inches(0.4),
                 font_size=15, bold=True, color=GREEN)

    # 할인율 공식
    add_text_box(sl, "할인율 = 5% + (위험점수 − 0.5) × 26%   (5% ≤ 할인율 ≤ 18%)",
                 Inches(0.6), Inches(6.4), Inches(11), Inches(0.35),
                 font_size=13, color=INK3, italic=True)

def slide_06(prs):
    """앱 데모"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 6)
    add_divider(sl)

    label_tag(sl, "앱 데모", Inches(0.6), Inches(0.68))
    add_text_box(sl, "예측이 결정으로 이어지는 두 화면",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)
    add_text_box(sl, "매니저가 바로 쓸 수 있는 Streamlit DSS 앱 — 취소 위험 리스트부터 Flexi 라우팅까지",
                 Inches(0.6), Inches(1.75), Inches(11.5), Inches(0.4),
                 font_size=14, color=INK2)

    tabs = [
        ("Tab 1", "예약 우선순위 리스트", [
            "취소 확률 내림차순으로 예약 정렬",
            "예약별 SHAP 위험 근거 3가지 자동 출력",
            '"왜 이 손님이 위험한가" 즉시 확인',
            "country · lead_time · special_requests",
        ], BRAND, BLUE_S),
        ("Tab 2", "Flexi 오버부킹 관리", [
            "오늘 Flexi 풀 현황 · 권장 슬롯 수 계산",
            "평균 할인율 · 예상 RevPAR 개선",
            "신규 예약 취소 위험 즉시 판별",
            "Flexi 라우팅 권장 + 매니저 승인 단계",
        ], ACCENT, GOLD_S),
    ]
    for i, (tag, ttl, bullets, c, bg_c) in enumerate(tabs):
        x = Inches(0.6 + i * 6.4)
        bg = sl.shapes.add_shape(1, x, Inches(2.35), Inches(6.1), Inches(6.1)
                                 if i == 0 else Inches(6.1))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_c
        bg.line.color.rgb = c
        bg.line.width = Pt(0.75)

        add_text_box(sl, tag, x+Inches(0.2), Inches(2.55), Inches(5.7), Inches(0.35),
                     font_size=11, bold=True, color=c)
        add_text_box(sl, ttl, x+Inches(0.2), Inches(2.9), Inches(5.7), Inches(0.5),
                     font_size=18, bold=True, color=INK)
        tb = sl.shapes.add_textbox(x+Inches(0.2), Inches(3.5), Inches(5.7), Inches(2.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"• {b}"
            run.font.size  = Pt(13)
            run.font.color.rgb = INK2

def slide_07(prs):
    """Q2 — 데이터 공백"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 7)
    add_divider(sl)

    label_tag(sl, "Q2 · 데이터 공백", Inches(0.6), Inches(0.68))
    add_text_box(sl, "취소가 예상보다 적게 발생하면 — 누구를 재배치하고, 얼마면 수락할까요?",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.85),
                 font_size=26, bold=True, color=INK)

    add_text_box(sl, "이 협상 이력은 데이터에 존재하지 않습니다.\n실제 호텔이 재배치 협상을 기록해야만 쌓이는 데이터입니다.",
                 Inches(0.6), Inches(2.1), Inches(11), Inches(0.7),
                 font_size=15, color=INK2)

    items = [
        ("✓", "취소 예측 데이터 있음",        True),
        ("✓", "Flexi 라우팅 설계 완료",       True),
        ("✗", "재배치 협상 이력 존재하지 않음", False),
    ]
    for i, (icon, text, ok) in enumerate(items):
        c = GREEN if ok else RED
        bg_c = RGBColor(0xF0,0xFD,0xF4) if ok else RGBColor(0xFF,0xF5,0xF5)
        bg = sl.shapes.add_shape(1, Inches(0.6), Inches(3.0 + i * 1.15),
                                 Inches(11.5), Inches(0.9))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_c
        bg.line.color.rgb = c
        bg.line.width = Pt(0.5)
        add_text_box(sl, f"{icon}  {text}",
                     Inches(0.85), Inches(3.05 + i * 1.15), Inches(11), Inches(0.55),
                     font_size=17, bold=ok, color=c if not ok else INK)

def slide_08(prs):
    """목표 모델 구조"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 8)
    add_divider(sl)

    label_tag(sl, "궁극의 목표 · 목표 모델 구조", Inches(0.6), Inches(0.68))
    add_text_box(sl, "목표 모델은 명확합니다 — 지금 데이터만 없습니다.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)

    add_text_box(sl, "모델이 학습할 데이터 — 현재 0행",
                 Inches(0.6), Inches(1.85), Inches(8), Inches(0.35),
                 font_size=12, bold=True, color=INK3)
    add_table(sl,
        ["lead_time", "country", "room", "meal", "오퍼 금액", "수락"],
        [
            ["45일", "GBR", "City",   "BB", "€46", "✓"],
            ["3일",  "FRA", "Resort", "HB", "€99", "✗"],
            ["21일", "DEU", "City",   "RO", "€66", "✓"],
        ],
        Inches(0.6), Inches(2.25), Inches(8.5), Inches(1.85),
        col_widths=[Inches(1.3), Inches(1.1), Inches(1.2), Inches(1.1), Inches(1.5), Inches(1.3)],
    )
    add_text_box(sl, "↑ 이런 행이 수천 건 쌓여야 학습 가능 — 지금은 한 행도 없음",
                 Inches(0.6), Inches(4.2), Inches(9), Inches(0.35),
                 font_size=12, italic=True, color=RED)

    add_text_box(sl, "학습 후 모델이 하는 일 (예시)",
                 Inches(0.6), Inches(4.7), Inches(9), Inches(0.35),
                 font_size=12, bold=True, color=INK3)
    add_text_box(sl, "입력: GBR · lead_time 45일 · City · BB\n→ 출력: 오퍼 €53 제안 시 수락 확률 87%",
                 Inches(0.6), Inches(5.1), Inches(9), Inches(0.6),
                 font_size=15, color=BRAND)

def slide_09(prs):
    """LLM 정당성 — Claim 0 → 1 → 2"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 9)
    add_divider(sl)

    label_tag(sl, "LLM 정당성 · Claim 0 → 1 → 2", Inches(0.6), Inches(0.68))
    add_text_box(sl, "데이터 공백을 논리로 메웁니다.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)

    # Claim 0 — 대안 비교
    alts = [
        ("ABM",       "✗", "행동 규칙 수작업 명세 필요 — 협상 도메인 전문가 의존, 신규 시나리오 취약"),
        ("설문/실험", "✗", "미출시 상품 반응 측정 불가 · 응답자 모집 비용·시간 과다"),
        ("통계 모델", "✗", "협상 이력 데이터 필요 — 그것이 바로 지금의 공백"),
        ("LLM",       "✓", "사전학습으로 호텔·협상 도메인 내재화 — 데이터 없이도 동작하는 유일한 합리적 선택"),
    ]
    add_text_box(sl, "Claim 0 · LLM이 유일하게 합리적인 대안입니다",
                 Inches(0.6), Inches(1.85), Inches(11), Inches(0.35),
                 font_size=12, bold=True, color=INK3)
    for i, (name, icon, desc) in enumerate(alts):
        ok = icon == "✓"
        c = GREEN if ok else RED
        bg_c = RGBColor(0xF0,0xFD,0xF4) if ok else RGBColor(0xFF,0xF5,0xF5)
        bg = sl.shapes.add_shape(1, Inches(0.6), Inches(2.25 + i * 0.95),
                                 Inches(12.2), Inches(0.78))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_c
        bg.line.color.rgb = c
        bg.line.width = Pt(0.5)
        add_text_box(sl, f"{icon} {name}",
                     Inches(0.8), Inches(2.3 + i * 0.95), Inches(1.5), Inches(0.4),
                     font_size=13, bold=True, color=c)
        add_text_box(sl, desc,
                     Inches(2.4), Inches(2.3 + i * 0.95), Inches(10), Inches(0.4),
                     font_size=13, color=INK2)

def slide_10(prs):
    """선행 연구 근거 — 3 논문"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 10)
    add_divider(sl)

    label_tag(sl, "선행 연구 근거 · 3 논문", Inches(0.6), Inches(0.68))
    add_text_box(sl, "세 논문이 시뮬레이션을 지지합니다.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)

    papers = [
        ("시뮬레이션",  "Park et al., UIST 2023",
         "LLM 에이전트가 인간의 사회적 행동·의사결정을 현실적으로 시뮬레이션함을 실증",
         "한계: 실험실 환경 — 실제 배포 검증 미비"),
        ("재현성",      "Aher et al., ICML 2023",
         "LLM이 인간 피험자 연구 결과를 통계적으로 재현. 합성 레이블 신뢰성 근거",
         "한계: 호텔 협상 특화 실험 아님"),
        ("방법론",      "Horton & Argyle, 2023",
         "LLM 시뮬레이션이 가격 협상 실험 결과를 재현. 우리 접근의 직접적 방법론 근거",
         "한계: 1회 협상 — 다회전 협상 검증 미비"),
    ]
    for i, (tag, ref, desc, limit) in enumerate(papers):
        x = Inches(0.6 + i * 4.25)
        bg = sl.shapes.add_shape(1, x, Inches(2.0), Inches(4.0), Inches(4.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLUE_S
        bg.line.color.rgb = BRAND
        bg.line.width = Pt(0.75)
        add_text_box(sl, tag, x+Inches(0.2), Inches(2.15), Inches(3.6), Inches(0.35),
                     font_size=11, bold=True, color=ACCENT)
        add_text_box(sl, ref, x+Inches(0.2), Inches(2.5),  Inches(3.6), Inches(0.4),
                     font_size=14, bold=True, color=BRAND)
        add_text_box(sl, desc, x+Inches(0.2), Inches(2.95), Inches(3.6), Inches(1.5),
                     font_size=13, color=INK2)
        add_text_box(sl, limit, x+Inches(0.2), Inches(4.55), Inches(3.6), Inches(0.7),
                     font_size=11, italic=True, color=INK3)

def slide_11(prs):
    """Claim 2 — 1,240 runs"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 11)
    add_divider(sl)

    label_tag(sl, "Claim 2 · 1,240 runs · hint 제거 최종본", Inches(0.6), Inches(0.68))
    add_text_box(sl, "유형마다 다른 수락 패턴이 나타났습니다.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)

    archetypes = [
        ("A", "비즈니스 솔로",    "재배치 용이",        BRAND),
        ("B", "레저 커플",        "앵커링 민감",        ACCENT),
        ("C", "가족",             "재배치 제외",        RED),
        ("D", "가격민감형 고객",  "임계값 반응형",      GREEN),
        ("E", "단체",             "별도 협상",          INK2),
    ]
    for i, (code, name, trait, color) in enumerate(archetypes):
        x = Inches(0.6 + i * 2.5)
        bg = sl.shapes.add_shape(1, x, Inches(2.0), Inches(2.3), Inches(2.2))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLUE_S
        bg.line.color.rgb = color
        bg.line.width = Pt(1.0)
        add_text_box(sl, code,  x+Inches(0.1), Inches(2.1),  Inches(0.5), Inches(0.5),
                     font_size=24, bold=True, color=color)
        add_text_box(sl, name,  x+Inches(0.1), Inches(2.65), Inches(2.1), Inches(0.4),
                     font_size=13, bold=True, color=INK)
        add_text_box(sl, trait, x+Inches(0.1), Inches(3.05), Inches(2.1), Inches(0.4),
                     font_size=12, color=INK2)

    # 결과 요약
    add_text_box(sl, "수락률 곡선 핵심 발견",
                 Inches(0.6), Inches(4.4), Inches(12), Inches(0.35),
                 font_size=12, bold=True, color=INK3)
    findings = [
        "A (비즈니스): 전 구간 100% — 가격 무관 수락",
        "D (예산형): €46 계단 임계 — step function 확인",
        "B (레저): €91 앵커링 역설 발견 — 높은 오퍼가 오히려 수락률 하락",
        "C, E: 전 구간 0% — 재배치 불가 유형 식별",
    ]
    for i, f in enumerate(findings):
        c = GREEN if "100%" in f else (RED if "역설" in f else INK2)
        add_text_box(sl, f"• {f}",
                     Inches(0.6), Inches(4.85 + i * 0.5), Inches(12), Inches(0.4),
                     font_size=13, color=c)

def slide_12(prs):
    """보완 실험 — 240 runs"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 12)
    add_divider(sl)

    label_tag(sl, "보완 실험 · 240 추가 runs 신규", Inches(0.6), Inches(0.68))
    add_text_box(sl, "임계값과 역설 — 그래프가 말하는 것.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)
    add_text_box(sl, "D 아키타입 · 임계값 정밀화 (room_cost = €110)",
                 Inches(0.6), Inches(1.8), Inches(8), Inches(0.35),
                 font_size=12, bold=True, color=INK3)

    add_table(sl,
        ["오퍼", "오퍼 비율", "수락률", "n", "비고"],
        [
            ["€29", "26%", "0%",    "40", ""],
            ["€33", "30%", "0%",    "40", "보완"],
            ["€37", "34%", "0%",    "40", "보완"],
            ["€42", "38%", "12.5%", "40", "보완"],
            ["€46", "42%", "100%",  "40", "임계값"],
            ["€66", "60%", "97.5%", "40", ""],
            ["€81", "74%", "100%",  "40", ""],
            ["€99", "90%", "100%",  "40", ""],
        ],
        Inches(0.6), Inches(2.2), Inches(7.5), Inches(3.6),
        col_widths=[Inches(1.0), Inches(1.3), Inches(1.3), Inches(0.8), Inches(1.3)],
    )

    add_text_box(sl, "Step function 확인: €42(12.5%) → €46(100%)\n→ 임계값은 €44~€46 구간",
                 Inches(0.6), Inches(5.9), Inches(8), Inches(0.6),
                 font_size=14, bold=True, color=BRAND)

    add_text_box(sl, "B(레저) 앵커링 역설",
                 Inches(8.8), Inches(2.2), Inches(4), Inches(0.35),
                 font_size=12, bold=True, color=INK3)
    add_text_box(sl, "€91 오퍼에서 수락률 하락\n높은 오퍼 → 오히려 불신 유발\n→ 협상 전략 수정 필요",
                 Inches(8.8), Inches(2.6), Inches(4.2), Inches(1.2),
                 font_size=13, color=RED)

def slide_13(prs):
    """궁극의 비전 — 미래 모델"""
    sl = blank_slide(prs)
    set_bg(sl, WHITE)
    add_header(sl, 13)
    add_divider(sl)

    label_tag(sl, "궁극의 비전 · 미래 모델", Inches(0.6), Inches(0.68))
    add_text_box(sl, "LLM은 데이터 공백의 브릿지입니다.",
                 Inches(0.6), Inches(1.05), Inches(12), Inches(0.65),
                 font_size=28, bold=True, color=INK)
    add_text_box(sl, "실 협상 데이터가 누적되면 ML이 이 역할을 대체합니다 — 우리 앱이 바로 그 데이터를 모읍니다.",
                 Inches(0.6), Inches(1.75), Inches(11.5), Inches(0.4),
                 font_size=14, color=INK2)

    stages = [
        ("1 지금",
         "캐글 + LLM 합성",
         "119K 예약으로 취소 패턴 학습\nLLM으로 재배치 협상 이력 합성\nFlexi 운영 가설 검증 + DSS 앱 배포",
         BLUE_S, BRAND),
        ("2 앱 운용 후",
         "실 데이터 + 하이브리드",
         "앱이 실 운영 데이터 수집 시작\n실제 협상 기록이 앱을 통해 누적\nLLM 예측 + 실데이터 병행",
         GOLD_S, ACCENT),
        ("3 궁극의 모델",
         "Pure ML 협상 모델",
         "실 협상 데이터로 지도학습\nLLM 의존도 제거\n실시간 개인화 오퍼 생성",
         RGBColor(0xF0,0xFD,0xF4), GREEN),
    ]
    for i, (phase, ttl, body, bg_c, c) in enumerate(stages):
        x = Inches(0.6 + i * 4.25)
        bg = sl.shapes.add_shape(1, x, Inches(2.3), Inches(4.0), Inches(4.4))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_c
        bg.line.color.rgb = c
        bg.line.width = Pt(0.75)
        add_text_box(sl, phase, x+Inches(0.2), Inches(2.45), Inches(3.6), Inches(0.35),
                     font_size=11, bold=True, color=c)
        add_text_box(sl, ttl,   x+Inches(0.2), Inches(2.85), Inches(3.6), Inches(0.5),
                     font_size=16, bold=True, color=INK)
        add_text_box(sl, body,  x+Inches(0.2), Inches(3.4),  Inches(3.6), Inches(2.5),
                     font_size=12, color=INK2)

def slide_14(prs):
    """Closing"""
    sl = blank_slide(prs)
    set_bg(sl, RGBColor(0x06, 0x09, 0x1A))  # 다크 배경
    add_header(sl, 14, dark=True)
    add_divider(sl, dark=True)

    add_text_box(sl, "Closing",
                 Inches(0.6), Inches(0.68), Inches(4), Inches(0.35),
                 font_size=10, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    add_text_box(sl, "저희는",
                 Inches(0.6), Inches(1.1), Inches(10), Inches(0.7),
                 font_size=48, bold=True, color=RGBColor(0xF1,0xF5,0xF9))
    add_text_box(sl, "두뇌를 만듭니다.",
                 Inches(0.6), Inches(1.8), Inches(10), Inches(0.7),
                 font_size=48, bold=True, color=ACCENT)

    results = [
        ("Q1", "LightGBM + 점유율 임계값 → RevPAR +81%"),
        ("Q2", "LLM 시뮬레이션 1,240 runs → 유형별 임계값 실측 + 앵커링 역설 발견"),
        ("앱", "Streamlit DSS — Tab 1 우선순위 · Tab 2 Flexi 관리 · Tab 3 오늘 예약 관리"),
    ]
    for i, (tag, desc) in enumerate(results):
        add_text_box(sl, tag,  Inches(0.6),      Inches(2.75 + i * 0.75),
                     Inches(0.8), Inches(0.5), font_size=14, bold=True, color=ACCENT)
        add_text_box(sl, desc, Inches(1.5),      Inches(2.75 + i * 0.75),
                     Inches(11), Inches(0.5),  font_size=14, color=RGBColor(0xF1,0xF5,0xF9))

    add_text_box(sl,
                 "PMS 연동 · 결제 시스템 · 협력 호텔 네트워크는 범위 밖입니다.\n저희는 취소 예측 엔진 + 재배치 보상 전략 두뇌를 만듭니다.",
                 Inches(0.6), Inches(5.2), Inches(11.5), Inches(0.8),
                 font_size=14, color=RGBColor(0x94,0xA3,0xB8))

    add_text_box(sl, "심재형 · 이고은 · 김나리  |  한양대학교  |  2026.06.10",
                 Inches(0.6), Inches(6.8), Inches(11.5), Inches(0.35),
                 font_size=11, color=RGBColor(0x47,0x55,0x69))

# ═══════════════════════════════════════════════════════════

def build_pptx(output_path: str):
    prs = new_prs()
    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    prs.save(output_path)
    print(f"Saved: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else r"c:\Users\jhsim\Erica261\M.L\projects\07_Hotel_DSS\presentations\Hotel_No-Show_DSS_v12.pptx"
    build_pptx(out)
